"""
routers/file_router.py
──────────────────────
Serves files from the local file system.
Extracts readable text from: PDF, DOCX, PPTX, XLSX, MD, TXT, images (OCR), video metadata.
All processing is done locally with open-source libraries only.
"""

from fastapi import APIRouter, UploadFile, File, HTTPException, Query
from fastapi.responses import FileResponse
from pathlib import Path
from typing import Optional
import os
import shutil
import mimetypes

from config import settings

router = APIRouter()

UPLOAD_DIR = Path(settings.UPLOAD_DIR).resolve()
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def safe_path(filename: str) -> Path:
    """Prevent directory traversal attacks."""
    target = (UPLOAD_DIR / filename).resolve()
    if not str(target).startswith(str(UPLOAD_DIR)):
        raise HTTPException(status_code=400, detail="Invalid file path")
    return target


def extract_text(filepath: Path) -> str:
    """
    Extract plain text from various file types.
    Returns extracted text, or a descriptive message if extraction is not possible.
    Never invents content — only returns what is genuinely in the file.
    """
    suffix = filepath.suffix.lower()
    
    if not suffix:
      try:
        return filepath.read_text(encoding="utf-8", errors="replace")
      except Exception as e:
        return f"[Read error: {e}]"
        
    # ── Plain text / Markdown ────────────────────────────────────────────────
    if suffix in (".txt", ".md", ".csv", ".log",""):
        try:
            return filepath.read_text(encoding="utf-8", errors="replace")
        except Exception as e:
            return f"[Text read error: {e}]"

    # ── PDF ──────────────────────────────────────────────────────────────────
    if suffix == ".pdf":
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
            return "\n".join(text_parts) if text_parts else "[PDF: no extractable text found]"
        except ImportError:
            return "[PDF extraction unavailable: install pdfplumber]"
        except Exception as e:
            return f"[PDF error: {e}]"

    # ── Word DOCX ────────────────────────────────────────────────────────────
    if suffix in (".docx", ".doc"):
        try:
            import docx
            doc = docx.Document(str(filepath))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return "[DOCX extraction unavailable: install python-docx]"
        except Exception as e:
            return f"[DOCX error: {e}]"

    # ── PowerPoint ───────────────────────────────────────────────────────────
    if suffix in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            parts = []
            for slide_no, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    parts.append(f"--- Slide {slide_no} ---\n" + "\n".join(slide_text))
            return "\n".join(parts) if parts else "[PPTX: no text found]"
        except ImportError:
            return "[PPTX extraction unavailable: install python-pptx]"
        except Exception as e:
            return f"[PPTX error: {e}]"

    # ── Excel / CSV ──────────────────────────────────────────────────────────
    if suffix in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts) if parts else "[XLSX: no data]"
        except ImportError:
            return "[XLSX extraction unavailable: install openpyxl]"
        except Exception as e:
            return f"[XLSX error: {e}]"

    # ── Images — OCR ─────────────────────────────────────────────────────────
    if suffix in (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff"):
        try:
            from PIL import Image
            import pytesseract
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            return text.strip() if text.strip() else "[Image: no text found via OCR]"
        except ImportError:
            return "[Image OCR unavailable: install Pillow and pytesseract + Tesseract binary]"
        except Exception as e:
            return f"[Image OCR error: {e}]"

    # ── Video — metadata only ─────────────────────────────────────────────────
    if suffix in (".mp4", ".avi", ".mov", ".mkv"):
        try:
            import subprocess, json as _json
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-print_format", "json",
                 "-show_format", "-show_streams", str(filepath)],
                capture_output=True, text=True, timeout=15
            )
            if result.returncode == 0:
                meta = _json.loads(result.stdout)
                fmt = meta.get("format", {})
                return (
                    f"Video file: {filepath.name}\n"
                    f"Duration: {fmt.get('duration', 'unknown')}s\n"
                    f"Size: {fmt.get('size', 'unknown')} bytes\n"
                    f"Format: {fmt.get('format_long_name', 'unknown')}"
                )
        except FileNotFoundError:
            return "[Video metadata unavailable: install ffprobe]"
        except Exception as e:
            return f"[Video metadata error: {e}]"

    return f"[Unsupported file type: {suffix}]"


# ─── Endpoints ────────────────────────────────────────────────────────────────

@router.get("/list")
def list_files(subdir: Optional[str] = Query(None)):
    """FastAPI endpoint — calls the internal function."""
    return list_files_internal(subdir)
def list_files_internal(subdir: Optional[str] = None) -> dict:
    """Pure Python function — callable from other modules without FastAPI Query."""
    base = UPLOAD_DIR if not subdir else (UPLOAD_DIR / subdir).resolve()
    if not base.exists():
        return {"file_count": 0, "files": []}

    files = []
    for f in sorted(base.rglob("*")):
        if f.is_file():
            files.append({
                "name": f.name,
                "relative_path": str(f.relative_to(UPLOAD_DIR)),
                "size_bytes": f.stat().st_size,
                "mime_type": mimetypes.guess_type(f.name)[0] or "application/octet-stream",
                "extension": f.suffix.lower(),
            })
    return {"file_count": len(files), "files": files}


@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload a file to local storage."""
    suffix = Path(file.filename).suffix.lower()
    if suffix not in settings.ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"File type '{suffix}' is not allowed. Allowed: {settings.ALLOWED_EXTENSIONS}"
        )

    dest = UPLOAD_DIR / file.filename
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)

    size_mb = dest.stat().st_size / (1024 * 1024)
    if size_mb > settings.MAX_FILE_SIZE_MB:
        dest.unlink()
        raise HTTPException(status_code=413, detail=f"File exceeds {settings.MAX_FILE_SIZE_MB} MB limit")

    return {"message": "Uploaded successfully", "filename": file.filename, "size_mb": round(size_mb, 2)}


@router.get("/read/{filename:path}")
def read_file(filename: str):
    """
    Extract and return text content from a file.
    For images: runs OCR. For PDFs/DOCX/PPTX: extracts text.
    Only real file content is returned — no generated text.
    """
    path = safe_path(filename)
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    content = extract_text(path)
    return {
        "filename": filename,
        "size_bytes": path.stat().st_size,
        "content": content,
        "chars": len(content),
    }


@router.get("/download/{filename:path}")
def download_file(filename: str):
    """Download a file directly."""
    path = safe_path(filename)
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(path, filename=path.name)


@router.delete("/delete/{filename:path}")
def delete_file(filename: str):
    """Delete a file from local storage."""
    path = safe_path(filename)
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    path.unlink()
    return {"message": f"Deleted: {filename}"}
